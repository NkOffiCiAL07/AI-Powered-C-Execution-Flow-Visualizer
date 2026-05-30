"""
Traceon Project — PowerPoint Generator
Run: python generate_ppt.py
Output: Traceon_Project.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Inches, Pt
import copy

# ─── Brand colours ──────────────────────────────────────────────────────────
ORANGE      = RGBColor(0xC9, 0x6A, 0x48)   # accent
ORANGE_DARK = RGBColor(0xA0, 0x4E, 0x30)   # darker accent
BG_DARK     = RGBColor(0x12, 0x12, 0x16)   # near-black slide bg
BG_MID      = RGBColor(0x1E, 0x1E, 0x2A)   # card bg
BG_LIGHT    = RGBColor(0xF5, 0xF0, 0xEB)   # off-white
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_LIGHT  = RGBColor(0xCC, 0xCC, 0xCC)
GRAY_MID    = RGBColor(0x88, 0x88, 0x99)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank

# ─── Helper utilities ────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=None, radius=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background() if line is None else None
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        if line_w:
            shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h, size=20, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_para(tf, text, size=16, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             space_before=None, bullet=False, level=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.level = level
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = ("• " if bullet else "") + text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def slide_bg(slide, color=BG_DARK):
    """Fill entire slide background."""
    add_rect(slide, 0, 0, W, H, fill=color)


def accent_bar(slide, top=False):
    """Thin orange horizontal accent stripe."""
    y = 0 if top else H - Inches(0.06)
    add_rect(slide, 0, y, W, Inches(0.06), fill=ORANGE)


def slide_number(slide, n, total):
    add_text(slide, f"{n} / {total}", W - Inches(1.2), H - Inches(0.45),
             Inches(1.1), Inches(0.35), size=10, color=GRAY_MID, align=PP_ALIGN.RIGHT)


# ─── SLIDE CONTENT ───────────────────────────────────────────────────────────
TOTAL = 13

# ══════════════════════════════════════════════════════════════════════════════
# 1 ▸ TITLE SLIDE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_bg(slide)

# Big orange diagonal accent block (top-right)
add_rect(slide, Inches(9.5), 0, Inches(3.83), Inches(3.5), fill=ORANGE)
# Darker triangle overlay effect (simple rect angled by color)
add_rect(slide, Inches(10.8), 0, Inches(2.53), Inches(2.0), fill=ORANGE_DARK)

# Logo badge circle (simulated with rect)
badge = add_rect(slide, Inches(0.55), Inches(2.0), Inches(1.1), Inches(1.1), fill=ORANGE)

add_text(slide, "T", Inches(0.62), Inches(2.05), Inches(1.0), Inches(1.0),
         size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Title
add_text(slide, "TRACEON", Inches(1.85), Inches(1.9), Inches(7.5), Inches(1.1),
         size=64, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# Tagline
add_text(slide, "AI-Powered Execution Flow Visualizer",
         Inches(1.87), Inches(3.05), Inches(7.2), Inches(0.7),
         size=24, color=ORANGE, bold=False)

# Sub-tagline
add_text(slide, "C  ·  C++  ·  Python  ·  Java",
         Inches(1.87), Inches(3.75), Inches(5.5), Inches(0.5),
         size=16, color=GRAY_LIGHT)

# Bottom bar
add_rect(slide, 0, H - Inches(1.0), W, Inches(1.0), fill=BG_MID)
add_text(slide, "React 18  ·  FastAPI  ·  LLDB  ·  Gemini AI  ·  Monaco Editor",
         Inches(0.5), H - Inches(0.75), Inches(9.0), Inches(0.5),
         size=13, color=GRAY_MID)
add_text(slide, "frontend-gamma-vert-20.vercel.app",
         Inches(9.5), H - Inches(0.75), Inches(3.5), Inches(0.5),
         size=12, color=ORANGE, align=PP_ALIGN.RIGHT)

accent_bar(slide, top=True)
slide_number(slide, 1, TOTAL)


# ══════════════════════════════════════════════════════════════════════════════
# 2 ▸ WHAT IS TRACEON?
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_bg(slide)
accent_bar(slide, top=True)
accent_bar(slide)

# Section header strip
add_rect(slide, 0, Inches(0.55), W, Inches(0.8), fill=BG_MID)
add_text(slide, "WHAT IS TRACEON?", Inches(0.5), Inches(0.6),
         Inches(8), Inches(0.65), size=28, bold=True, color=ORANGE)

# Left column — description
tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.55), Inches(6.0), Inches(4.8))
tf = tb.text_frame
tf.word_wrap = True

add_para(tf, "Traceon is a full-stack development tool that makes program execution visible — step by step, line by line.",
         size=17, color=WHITE, space_before=0)
add_para(tf, "", size=6, color=WHITE)
add_para(tf, "Developers use it to:", size=14, color=GRAY_LIGHT, space_before=4)
add_para(tf, "Understand how code actually runs at runtime", size=15, color=WHITE, bullet=True)
add_para(tf, "Catch bugs by watching variables change live", size=15, color=WHITE, bullet=True)
add_para(tf, "Visualise function call graphs automatically", size=15, color=WHITE, bullet=True)
add_para(tf, "Get AI-generated explanations and optimisations", size=15, color=WHITE, bullet=True)
add_para(tf, "Share code snippets via encoded URL", size=15, color=WHITE, bullet=True)
add_para(tf, "", size=6, color=WHITE)
add_para(tf, "No setup. No plugin. Just paste your code and debug.", size=14, color=ORANGE, bold=True)

# Right column — stats cards
def stat_card(slide, l, t, w, h, number, label):
    add_rect(slide, l, t, w, h, fill=BG_MID, line=ORANGE, line_w=Pt(1.5))
    add_text(slide, number, l, t + Inches(0.12), w, Inches(0.6),
             size=32, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    add_text(slide, label, l, t + Inches(0.65), w, Inches(0.45),
             size=12, color=GRAY_LIGHT, align=PP_ALIGN.CENTER)

stat_card(slide, Inches(7.0),  Inches(1.6),  Inches(2.7), Inches(1.2), "4",       "Languages")
stat_card(slide, Inches(10.1), Inches(1.6),  Inches(2.7), Inches(1.2), "10+",     "Major Features")
stat_card(slide, Inches(7.0),  Inches(3.0),  Inches(2.7), Inches(1.2), "5",       "Editor Themes")
stat_card(slide, Inches(10.1), Inches(3.0),  Inches(2.7), Inches(1.2), "Gemini",  "AI Engine")
stat_card(slide, Inches(7.0),  Inches(4.4),  Inches(2.7), Inches(1.2), "LLDB",    "Debugger Core")
stat_card(slide, Inches(10.1), Inches(4.4),  Inches(2.7), Inches(1.2), "Vercel",  "Deployment")

slide_number(slide, 2, TOTAL)


# ══════════════════════════════════════════════════════════════════════════════
# 3 ▸ SYSTEM ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_bg(slide)
accent_bar(slide, top=True)
accent_bar(slide)

add_rect(slide, 0, Inches(0.55), W, Inches(0.8), fill=BG_MID)
add_text(slide, "SYSTEM ARCHITECTURE", Inches(0.5), Inches(0.6),
         Inches(10), Inches(0.65), size=28, bold=True, color=ORANGE)

# Three columns
cols = [
    ("FRONTEND", Inches(0.35), [
        "React 18",
        "Monaco Editor",
        "CSS Variables",
        "5 Themes",
        "State-based routing",
        "Google OAuth (popup)",
        "JWT localStorage",
        "Vercel Deploy",
    ]),
    ("BACKEND", Inches(4.9), [
        "FastAPI + Uvicorn",
        "LLDB (C / C++)",
        "sys.settrace (Python)",
        "JVM Tracer (Java)",
        "Gemini AI API",
        "Google OAuth2",
        "JTI Blocklist",
        "MongoDB (optional)",
    ]),
    ("AI LAYER", Inches(9.4), [
        "Gemini 1.5 Flash",
        "Code Generation",
        "Execution Explanation",
        "Performance Optimiser",
        "Hotspot Analysis",
        "Prompt Engineering",
        "Streaming Responses",
        "Rate Limiting",
    ]),
]

for title, lx, items in cols:
    add_rect(slide, lx, Inches(1.55), Inches(4.0), Inches(5.55), fill=BG_MID, line=ORANGE, line_w=Pt(1))
    add_rect(slide, lx, Inches(1.55), Inches(4.0), Inches(0.55), fill=ORANGE)
    add_text(slide, title, lx + Inches(0.1), Inches(1.6), Inches(3.8), Inches(0.45),
             size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb = slide.shapes.add_textbox(lx + Inches(0.2), Inches(2.2), Inches(3.6), Inches(4.7))
    tf = tb.text_frame
    tf.word_wrap = True
    for item in items:
        add_para(tf, item, size=14, color=GRAY_LIGHT, bullet=True)

# Arrows between columns (simple text)
add_text(slide, "⟷", Inches(4.45), Inches(4.0), Inches(0.5), Inches(0.5),
         size=22, color=ORANGE, align=PP_ALIGN.CENTER)
add_text(slide, "⟷", Inches(8.95), Inches(4.0), Inches(0.5), Inches(0.5),
         size=22, color=ORANGE, align=PP_ALIGN.CENTER)

add_text(slide, "Port 3000", Inches(0.35), Inches(7.1), Inches(4.0), Inches(0.3),
         size=11, color=GRAY_MID, align=PP_ALIGN.CENTER)
add_text(slide, "Port 8000", Inches(4.9), Inches(7.1), Inches(4.0), Inches(0.3),
         size=11, color=GRAY_MID, align=PP_ALIGN.CENTER)
add_text(slide, "Gemini API", Inches(9.4), Inches(7.1), Inches(4.0), Inches(0.3),
         size=11, color=GRAY_MID, align=PP_ALIGN.CENTER)

slide_number(slide, 3, TOTAL)


# ══════════════════════════════════════════════════════════════════════════════
# 4 ▸ KEY FEATURES — EDITOR
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_bg(slide)
accent_bar(slide, top=True)
accent_bar(slide)

add_rect(slide, 0, Inches(0.55), W, Inches(0.8), fill=BG_MID)
add_text(slide, "SMART CODE EDITOR", Inches(0.5), Inches(0.6),
         Inches(10), Inches(0.65), size=28, bold=True, color=ORANGE)

features = [
    ("Monaco Editor", "VS Code-grade editor with syntax highlighting, auto-complete, and IntelliSense-style UX across all 4 languages."),
    ("5 Editor Themes", "Light · Dark · Ocean · Forest · Midnight — CSS-variable-driven, zero flicker on switch."),
    ("Live Syntax Check", "Clang -fsyntax-only runs debounced every 1500 ms. Errors appear as red squiggles inline."),
    ("Breakpoint Gutter", "Click any line number to toggle a red breakpoint dot. Persisted to localStorage. Auto-pauses during playback."),
    ("Execution Heatmap", "After a debug run, gutter bars colour each line blue→yellow→red based on hit count. Shows hot paths instantly."),
    ("Share Code Button", "Encodes code + language into a URL hash (base64). One click copies the link — shareable with no backend."),
]

cols_per_row = 3
for i, (title, desc) in enumerate(features):
    col = i % cols_per_row
    row = i // cols_per_row
    lx = Inches(0.35) + col * Inches(4.32)
    ty = Inches(1.6)  + row * Inches(2.65)
    add_rect(slide, lx, ty, Inches(4.1), Inches(2.45), fill=BG_MID, line=ORANGE, line_w=Pt(1))
    add_rect(slide, lx, ty, Inches(4.1), Inches(0.5), fill=ORANGE)
    add_text(slide, title, lx + Inches(0.12), ty + Inches(0.05), Inches(3.86), Inches(0.42),
             size=14, bold=True, color=WHITE)
    add_text(slide, desc, lx + Inches(0.15), ty + Inches(0.6), Inches(3.8), Inches(1.75),
             size=12, color=GRAY_LIGHT)

slide_number(slide, 4, TOTAL)


# ══════════════════════════════════════════════════════════════════════════════
# 5 ▸ KEY FEATURES — DEBUGGER
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_bg(slide)
accent_bar(slide, top=True)
accent_bar(slide)

add_rect(slide, 0, Inches(0.55), W, Inches(0.8), fill=BG_MID)
add_text(slide, "EXECUTION DEBUGGER & FLOW VISUALIZER", Inches(0.5), Inches(0.6),
         Inches(11), Inches(0.65), size=28, bold=True, color=ORANGE)

dbg_features = [
    ("Step Controls",       "Step Over · Step Into · Step Out · Back — full bi-directional navigation of execution state."),
    ("Auto-Play",           "Configurable speed (150 ms–2000 ms). Toggle Step In vs Step Over mode during playback."),
    ("Timeline Scrubber",   "Drag or click the progress bar to jump to any visited step. Red tick marks show breakpoint positions."),
    ("Breakpoint Auto-Pause","Execution auto-halts when a breakpointed line is hit during play. Shows a Resume badge."),
    ("Performance Badge",   "⚡ Blazing (≤20 steps) · 🟡 Moderate (≤120) · 🔴 Heavy (>120) — shown next to the step counter."),
    ("Call Graph SVG",      "Live SVG showing function call relationships built from the execution trace."),
    ("Variable Change Cards","Every step shows which variables changed, their old/new values, and the current call stack."),
    ("Keyboard Shortcuts",  "← → arrow keys, Space (play/pause), F5 (restart) — no mouse needed for stepping."),
]

for i, (title, desc) in enumerate(dbg_features):
    col = i % cols_per_row
    row = i // cols_per_row
    lx = Inches(0.35) + col * Inches(4.32)
    ty = Inches(1.6)  + row * Inches(2.65)
    if i == 6:
        lx = Inches(0.35) + 0 * Inches(4.32)
        ty = Inches(1.6) + 2 * Inches(2.65)
    if i == 7:
        lx = Inches(0.35) + 1 * Inches(4.32)
        ty = Inches(1.6) + 2 * Inches(2.65)
    add_rect(slide, lx, ty, Inches(4.1), Inches(2.45), fill=BG_MID, line=ORANGE, line_w=Pt(1))
    add_rect(slide, lx, ty, Inches(4.1), Inches(0.5), fill=ORANGE)
    add_text(slide, title, lx + Inches(0.12), ty + Inches(0.05), Inches(3.86), Inches(0.42),
             size=14, bold=True, color=WHITE)
    add_text(slide, desc, lx + Inches(0.15), ty + Inches(0.6), Inches(3.8), Inches(1.75),
             size=12, color=GRAY_LIGHT)

slide_number(slide, 5, TOTAL)


# ══════════════════════════════════════════════════════════════════════════════
# 6 ▸ CODE FLOW GRAPH
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_bg(slide)
accent_bar(slide, top=True)
accent_bar(slide)

add_rect(slide, 0, Inches(0.55), W, Inches(0.8), fill=BG_MID)
add_text(slide, "CODE FLOW GRAPH", Inches(0.5), Inches(0.6),
         Inches(10), Inches(0.65), size=28, bold=True, color=ORANGE)

# Left: description
tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.55), Inches(6.0), Inches(5.5))
tf = tb.text_frame
tf.word_wrap = True
add_para(tf, "Interactive SVG flowchart generated from the live execution trace.", size=16, color=WHITE, bold=True)
add_para(tf, "", size=5, color=WHITE)
add_para(tf, "Node Types", size=14, color=ORANGE, bold=True, space_before=4)
for node in ["entry / exit", "condition (if/else)", "loop (for/while)", "call / return",
             "func-def", "noise (filtered out)"]:
    add_para(tf, node, size=13, color=GRAY_LIGHT, bullet=True)
add_para(tf, "", size=5, color=WHITE)
add_para(tf, "Interactions", size=14, color=ORANGE, bold=True, space_before=4)
for feat in ["Pan (drag), zoom (scroll wheel), reset", "Click any node for detail panel",
             "Detail shows: code, type, exec count, variables, call stack",
             "Back-edge detection for loop arrows", "Per-node clipPath (no text overflow)",
             "Resets automatically on new session"]:
    add_para(tf, feat, size=13, color=GRAY_LIGHT, bullet=True)

# Right: mock flow graph diagram
boxes = [
    (Inches(7.5),  Inches(1.6),  Inches(5.0), Inches(0.55), ORANGE,   "ENTRY: main()"),
    (Inches(7.5),  Inches(2.35), Inches(5.0), Inches(0.55), BG_MID,   "CONDITION: if (n > 0)"),
    (Inches(6.8),  Inches(3.1),  Inches(2.3), Inches(0.55), BG_MID,   "CALL: add(a, b)"),
    (Inches(9.5),  Inches(3.1),  Inches(2.7), Inches(0.55), BG_MID,   "LOOP: for i in range"),
    (Inches(6.8),  Inches(3.85), Inches(2.3), Inches(0.55), BG_MID,   "RETURN: result"),
    (Inches(9.5),  Inches(3.85), Inches(2.7), Inches(0.55), BG_MID,   "LOOP BODY: ×3"),
    (Inches(7.5),  Inches(4.6),  Inches(5.0), Inches(0.55), BG_MID,   "EXIT: return 0"),
    (Inches(7.5),  Inches(5.35), Inches(5.0), Inches(0.55), ORANGE_DARK, "END"),
]
for lx, ty, bw, bh, col, label in boxes:
    add_rect(slide, lx, ty, bw, bh, fill=col, line=ORANGE, line_w=Pt(1))
    add_text(slide, label, lx + Inches(0.1), ty + Inches(0.05), bw - Inches(0.2), bh - Inches(0.1),
             size=11, color=WHITE, align=PP_ALIGN.CENTER)

# Arrow indicators (unicode)
for ay in [Inches(2.2), Inches(2.95), Inches(3.7), Inches(4.45), Inches(5.2)]:
    add_text(slide, "↓", Inches(9.75), ay, Inches(0.5), Inches(0.4),
             size=14, color=ORANGE, align=PP_ALIGN.CENTER)

slide_number(slide, 6, TOTAL)


# ══════════════════════════════════════════════════════════════════════════════
# 7 ▸ AI INTEGRATION
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_bg(slide)
accent_bar(slide, top=True)
accent_bar(slide)

add_rect(slide, 0, Inches(0.55), W, Inches(0.8), fill=BG_MID)
add_text(slide, "AI INTEGRATION — GEMINI", Inches(0.5), Inches(0.6),
         Inches(10), Inches(0.65), size=28, bold=True, color=ORANGE)

ai_cards = [
    ("✦ Generate Code",
     "Natural language → working code.\nDescribe what you want, Gemini writes it directly into the Monaco editor.",
     ["Supports all 4 languages", "Fills editor on generation", "Loading state with spinner"]),
    ("✦ Explain Code",
     "Full AI walkthrough of what the current code does: logic, complexity (Big-O), edge cases, and suggestions.",
     ["Triggered via floating FAB", "Shows complexity analysis", "Markdown-rendered output"]),
    ("✦ Optimise Code",
     "Sends execution hotspot data + code to Gemini. Returns targeted performance recommendations with reasoning.",
     ["Uses heatmap data as context", "Identifies bottleneck lines", "Actionable suggestions"]),
]

for i, (title, desc, points) in enumerate(ai_cards):
    lx = Inches(0.35) + i * Inches(4.32)
    ty = Inches(1.55)
    add_rect(slide, lx, ty, Inches(4.1), Inches(5.55), fill=BG_MID, line=ORANGE, line_w=Pt(1.5))
    add_rect(slide, lx, ty, Inches(4.1), Inches(0.6), fill=ORANGE)
    add_text(slide, title, lx + Inches(0.15), ty + Inches(0.07), Inches(3.8), Inches(0.5),
             size=16, bold=True, color=WHITE)
    add_text(slide, desc, lx + Inches(0.2), ty + Inches(0.75), Inches(3.7), Inches(1.8),
             size=13, color=GRAY_LIGHT)
    tb = slide.shapes.add_textbox(lx + Inches(0.2), ty + Inches(2.7), Inches(3.7), Inches(2.7))
    tf = tb.text_frame
    tf.word_wrap = True
    add_para(tf, "Highlights", size=12, color=ORANGE, bold=True)
    for pt in points:
        add_para(tf, pt, size=12, color=WHITE, bullet=True)

# Bottom note
add_text(slide, "⚡  All AI features use the floating FAB pattern — always accessible bottom-right, never buried in toolbar menus.",
         Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.38),
         size=12, color=ORANGE, align=PP_ALIGN.CENTER, italic=True)

slide_number(slide, 7, TOTAL)


# ══════════════════════════════════════════════════════════════════════════════
# 8 ▸ LANGUAGE SUPPORT
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_bg(slide)
accent_bar(slide, top=True)
accent_bar(slide)

add_rect(slide, 0, Inches(0.55), W, Inches(0.8), fill=BG_MID)
add_text(slide, "LANGUAGE SUPPORT", Inches(0.5), Inches(0.6),
         Inches(10), Inches(0.65), size=28, bold=True, color=ORANGE)

# Table header
headers = ["Language", "Editor", "Run", "Debugger", "Live Check", "AI", "Flow Graph"]
col_widths = [Inches(2.0), Inches(1.5), Inches(1.3), Inches(2.5), Inches(2.0), Inches(1.5), Inches(2.0)]
start_x = Inches(0.35)
start_y = Inches(1.6)

# Header row
x = start_x
for j, (hdr, cw) in enumerate(zip(headers, col_widths)):
    add_rect(slide, x, start_y, cw, Inches(0.55), fill=ORANGE)
    add_text(slide, hdr, x + Inches(0.05), start_y + Inches(0.06), cw - Inches(0.1), Inches(0.45),
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    x += cw

# Data rows
rows = [
    ("C",      "✅", "✅", "✅ LLDB",        "✅ clang", "✅", "✅"),
    ("C++",    "✅", "✅", "✅ LLDB",        "✅ clang", "✅", "✅"),
    ("Python", "✅", "✅", "✅ sys.settrace","✅",       "✅", "✅"),
    ("Java",   "✅", "✅", "✅ JVM tracer",  "✅",       "✅", "✅"),
]

for ri, row in enumerate(rows):
    y = start_y + Inches(0.55) + ri * Inches(0.75)
    x = start_x
    row_fill = BG_MID if ri % 2 == 0 else RGBColor(0x18, 0x18, 0x24)
    for j, (cell, cw) in enumerate(zip(row, col_widths)):
        add_rect(slide, x, y, cw, Inches(0.72), fill=row_fill, line=RGBColor(0x33, 0x33, 0x44), line_w=Pt(0.5))
        txt_color = ORANGE if j == 0 else WHITE
        add_text(slide, cell, x + Inches(0.05), y + Inches(0.1), cw - Inches(0.1), Inches(0.55),
                 size=13, color=txt_color, align=PP_ALIGN.CENTER)
        x += cw

# Debugger detail boxes below table
detail_y = Inches(5.1)
details = [
    ("C / C++ — LLDB",    "Full LLDB session via Python API. Step into stdlib. View raw pointer values, struct members, and memory addresses."),
    ("Python",            "subprocess sys.settrace tracer. Captures local scope at every line. JSON snapshot output. No GIL issues."),
    ("Java",              "JVM-level tracer. Method entry/exit events, local variable inspection. Outputs compatible snapshot format."),
]
for i, (title, desc) in enumerate(details):
    lx = Inches(0.35) + i * Inches(4.32)
    add_rect(slide, lx, detail_y, Inches(4.1), Inches(1.9), fill=BG_MID, line=ORANGE, line_w=Pt(1))
    add_text(slide, title, lx + Inches(0.15), detail_y + Inches(0.1), Inches(3.8), Inches(0.45),
             size=13, bold=True, color=ORANGE)
    add_text(slide, desc, lx + Inches(0.15), detail_y + Inches(0.55), Inches(3.8), Inches(1.25),
             size=11, color=GRAY_LIGHT)

slide_number(slide, 8, TOTAL)


# ══════════════════════════════════════════════════════════════════════════════
# 9 ▸ AUTH & SECURITY
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_bg(slide)
accent_bar(slide, top=True)
accent_bar(slide)

add_rect(slide, 0, Inches(0.55), W, Inches(0.8), fill=BG_MID)
add_text(slide, "AUTH & SECURITY", Inches(0.5), Inches(0.6),
         Inches(10), Inches(0.65), size=28, bold=True, color=ORANGE)

# Flow diagram
flow_items = [
    ("User clicks\nSign In", Inches(0.5)),
    ("Google OAuth\nPopup", Inches(3.0)),
    ("Callback →\nFastAPI", Inches(5.5)),
    ("JWT Issued\n(HS256)", Inches(8.0)),
    ("localStorage\nToken", Inches(10.5)),
]
flow_y = Inches(1.9)
for label, lx in flow_items:
    add_rect(slide, lx, flow_y, Inches(2.2), Inches(0.9), fill=BG_MID, line=ORANGE, line_w=Pt(1.5))
    add_text(slide, label, lx + Inches(0.1), flow_y + Inches(0.05), Inches(2.0), Inches(0.8),
             size=12, color=WHITE, align=PP_ALIGN.CENTER)

for lx in [Inches(2.75), Inches(5.25), Inches(7.75), Inches(10.25)]:
    add_text(slide, "→", lx, flow_y + Inches(0.2), Inches(0.4), Inches(0.5),
             size=18, color=ORANGE, align=PP_ALIGN.CENTER)

# Security features
sec_cols = [
    ("JWT Token Security", [
        "HS256 signed tokens",
        "JTI (JWT ID) per token",
        "Server-side JTI blocklist",
        "Logout invalidates token immediately",
        "Configurable TTL (default 7 days)",
        "Token restored on page reload",
    ]),
    ("OAuth2 Flow", [
        "Google OAuth2 popup (no redirect)",
        "CSRF state parameter",
        "postMessage closes popup on success",
        "Redirect URI server-side only",
        "Rate-limited /auth/google endpoint",
        "AuthContext React context",
    ]),
    ("Guest Mode", [
        "provider: 'guest', role: 'guest'",
        "Only Run is available (no Debug/AI)",
        "Sign-in prompt on locked features",
        "Lock icon on AI FAB for guests",
        "No token stored",
        "Session expires on page close",
    ]),
]

for i, (title, items) in enumerate(sec_cols):
    lx = Inches(0.35) + i * Inches(4.32)
    ty = Inches(3.1)
    add_rect(slide, lx, ty, Inches(4.1), Inches(4.0), fill=BG_MID, line=ORANGE, line_w=Pt(1))
    add_rect(slide, lx, ty, Inches(4.1), Inches(0.5), fill=ORANGE)
    add_text(slide, title, lx + Inches(0.12), ty + Inches(0.06), Inches(3.86), Inches(0.42),
             size=14, bold=True, color=WHITE)
    tb = slide.shapes.add_textbox(lx + Inches(0.2), ty + Inches(0.6), Inches(3.7), Inches(3.3))
    tf = tb.text_frame
    tf.word_wrap = True
    for item in items:
        add_para(tf, item, size=13, color=GRAY_LIGHT, bullet=True)

slide_number(slide, 9, TOTAL)


# ══════════════════════════════════════════════════════════════════════════════
# 10 ▸ UI / UX HIGHLIGHTS
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_bg(slide)
accent_bar(slide, top=True)
accent_bar(slide)

add_rect(slide, 0, Inches(0.55), W, Inches(0.8), fill=BG_MID)
add_text(slide, "UI / UX HIGHLIGHTS", Inches(0.5), Inches(0.6),
         Inches(10), Inches(0.65), size=28, bold=True, color=ORANGE)

ux_items = [
    ("Smooth View Transitions",
     "Every page switch triggers a fade+slide keyframe animation. Zero layout jumps."),
    ("Onboarding Tour",
     "4-step spotlight tour for first-time users: targets real DOM elements, localStorage-gated."),
    ("Animated Stats Counter",
     "Landing page stats count up when scrolled into view via IntersectionObserver."),
    ("Testimonials Carousel",
     "Auto-rotates every 4 seconds, dot navigation, smooth swap animation."),
    ("Floating AI FAB",
     "Always-visible AI button, bottom-right. Pulses when ready, spinner when loading, lock when guest."),
    ("Performance Score Badge",
     "⚡ Blazing / 🟡 Moderate / 🔴 Heavy badge animates in next to the step counter."),
    ("Share Code Button",
     "Encodes code+language to URL hash. One click. No backend. Toast confirmation."),
    ("Responsive Theme System",
     "5 themes. CSS variables only. Monaco editor themes matched. Zero hardcoded colors."),
    ("Docs Page",
     "Full documentation: Flow Graph, Heatmap, Breakpoints, Share Code, AI Insights, Quick Start."),
]

for i, (title, desc) in enumerate(ux_items):
    col = i % 3
    row = i // 3
    lx = Inches(0.35) + col * Inches(4.32)
    ty = Inches(1.6)  + row * Inches(1.85)
    add_rect(slide, lx, ty, Inches(4.1), Inches(1.7), fill=BG_MID, line=ORANGE, line_w=Pt(0.75))
    add_text(slide, title, lx + Inches(0.15), ty + Inches(0.1), Inches(3.8), Inches(0.45),
             size=13, bold=True, color=ORANGE)
    add_text(slide, desc, lx + Inches(0.15), ty + Inches(0.55), Inches(3.8), Inches(1.05),
             size=11, color=GRAY_LIGHT)

slide_number(slide, 10, TOTAL)


# ══════════════════════════════════════════════════════════════════════════════
# 11 ▸ TECH STACK
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_bg(slide)
accent_bar(slide, top=True)
accent_bar(slide)

add_rect(slide, 0, Inches(0.55), W, Inches(0.8), fill=BG_MID)
add_text(slide, "FULL TECH STACK", Inches(0.5), Inches(0.6),
         Inches(10), Inches(0.65), size=28, bold=True, color=ORANGE)

stack_cols = [
    ("Frontend", [
        ("React 18",                "UI framework"),
        ("Monaco Editor",           "Code editor engine (VS Code core)"),
        ("pptxgenjs",               "PPTX generation"),
        ("@react-pdf/renderer",     "PDF generation"),
        ("CSS Variables",           "Theme system, no preprocessor"),
        ("Google Fonts",            "Inter + Material Symbols"),
        ("IntersectionObserver",    "Scroll-triggered animations"),
        ("localStorage",            "Token, breakpoints, tour flag"),
    ]),
    ("Backend", [
        ("FastAPI",                 "Async Python web framework"),
        ("Uvicorn",                 "ASGI server"),
        ("LLDB (Python API)",       "C/C++ debugger"),
        ("sys.settrace",            "Python execution tracer"),
        ("Gemini AI (google-generativeai)", "LLM API"),
        ("python-jose",             "JWT (HS256)"),
        ("httpx",                   "OAuth HTTP client"),
        ("reportlab / fpdf2",       "Server-side PDF generation"),
        ("python-pptx",             "Server-side PPTX generation"),
        ("MongoDB (optional)",      "Session persistence"),
    ]),
    ("Infrastructure", [
        ("Vercel",                  "Frontend hosting + CDN"),
        ("GitHub Actions",          "CI/CD deploy on push"),
        ("Google Cloud Console",    "OAuth2 credentials"),
        ("clang",                   "Live syntax checking"),
        ("Node.js 18+",             "Frontend build runtime"),
        ("Python 3.10+",            "Backend runtime"),
        ("JWT (HS256)",             "Stateless auth tokens"),
        ("CORS middleware",         "Origin whitelist"),
    ]),
]

for i, (title, items) in enumerate(stack_cols):
    lx = Inches(0.35) + i * Inches(4.32)
    ty = Inches(1.55)
    add_rect(slide, lx, ty, Inches(4.1), Inches(5.55), fill=BG_MID, line=ORANGE, line_w=Pt(1))
    add_rect(slide, lx, ty, Inches(4.1), Inches(0.5), fill=ORANGE)
    add_text(slide, title, lx + Inches(0.12), ty + Inches(0.05), Inches(3.86), Inches(0.42),
             size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb = slide.shapes.add_textbox(lx + Inches(0.15), ty + Inches(0.6), Inches(3.8), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    for tech, note in items:
        p = tf.add_paragraph()
        run1 = p.add_run()
        run1.text = "• " + tech
        run1.font.size = Pt(12)
        run1.font.bold = True
        run1.font.color.rgb = WHITE
        run2 = p.add_run()
        run2.text = " — " + note
        run2.font.size = Pt(11)
        run2.font.color.rgb = GRAY_MID

slide_number(slide, 11, TOTAL)


# ══════════════════════════════════════════════════════════════════════════════
# 12 ▸ DEPLOYMENT
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_bg(slide)
accent_bar(slide, top=True)
accent_bar(slide)

add_rect(slide, 0, Inches(0.55), W, Inches(0.8), fill=BG_MID)
add_text(slide, "DEPLOYMENT PIPELINE", Inches(0.5), Inches(0.6),
         Inches(10), Inches(0.65), size=28, bold=True, color=ORANGE)

# Pipeline steps
steps = [
    ("1", "Code Push", "Developer pushes to\nmain branch on GitHub"),
    ("2", "GH Actions", ".github/workflows/\ndeploy-frontend.yml fires"),
    ("3", "Vercel CLI", "vercel --prod deploys\nfrontend/ directory"),
    ("4", "SPA Rewrite", "vercel.json rewrites\nall routes → index.html"),
    ("5", "Live ✅",     "frontend-gamma-vert-20\n.vercel.app updated"),
]
sy = Inches(2.0)
for i, (num, title, desc) in enumerate(steps):
    lx = Inches(0.5) + i * Inches(2.55)
    add_rect(slide, lx, sy, Inches(2.2), Inches(2.3), fill=BG_MID, line=ORANGE, line_w=Pt(1.5))
    add_rect(slide, lx, sy, Inches(2.2), Inches(0.6), fill=ORANGE)
    add_text(slide, num, lx, sy + Inches(0.05), Inches(2.2), Inches(0.52),
             size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, lx + Inches(0.1), sy + Inches(0.72), Inches(2.0), Inches(0.45),
             size=13, bold=True, color=ORANGE)
    add_text(slide, desc, lx + Inches(0.1), sy + Inches(1.2), Inches(2.0), Inches(1.0),
             size=11, color=GRAY_LIGHT)
    if i < 4:
        add_text(slide, "→", lx + Inches(2.2), sy + Inches(0.85), Inches(0.35), Inches(0.5),
                 size=18, color=ORANGE, align=PP_ALIGN.CENTER)

# Secrets table
add_text(slide, "Required GitHub Secrets", Inches(0.5), Inches(4.65), Inches(6.0), Inches(0.45),
         size=15, bold=True, color=ORANGE)
secrets = [("VERCEL_TOKEN", "Personal access token"), ("VERCEL_ORG_ID", "Vercel team/org ID"),
           ("VERCEL_PROJECT_ID", "Linked project ID")]
for i, (k, v) in enumerate(secrets):
    ty = Inches(5.15) + i * Inches(0.55)
    add_rect(slide, Inches(0.5), ty, Inches(2.5), Inches(0.48), fill=BG_MID, line=ORANGE, line_w=Pt(0.5))
    add_text(slide, k, Inches(0.6), ty + Inches(0.05), Inches(2.3), Inches(0.38),
             size=12, color=ORANGE, bold=True)
    add_rect(slide, Inches(3.1), ty, Inches(3.0), Inches(0.48), fill=BG_MID, line=RGBColor(0x33,0x33,0x44), line_w=Pt(0.5))
    add_text(slide, v, Inches(3.2), ty + Inches(0.05), Inches(2.8), Inches(0.38),
             size=12, color=GRAY_LIGHT)

# Env vars
add_text(slide, "Key Environment Variables (.env)", Inches(6.8), Inches(4.65), Inches(6.0), Inches(0.45),
         size=15, bold=True, color=ORANGE)
envs = [
    "GEMINI_API_KEY — AI features",
    "GOOGLE_CLIENT_ID / SECRET — OAuth",
    "JWT_SECRET — token signing",
    "OAUTH_REDIRECT_URI — callback URL",
    "MONGO_URI — optional persistence",
]
tb = slide.shapes.add_textbox(Inches(6.8), Inches(5.15), Inches(6.0), Inches(2.2))
tf = tb.text_frame
tf.word_wrap = True
for env in envs:
    add_para(tf, env, size=12, color=GRAY_LIGHT, bullet=True)

slide_number(slide, 12, TOTAL)


# ══════════════════════════════════════════════════════════════════════════════
# 13 ▸ CLOSING / CTA
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_bg(slide)

# Full-width orange accent at top
add_rect(slide, 0, 0, W, Inches(0.12), fill=ORANGE)

# Background decorative block
add_rect(slide, Inches(9.0), 0, Inches(4.33), H, fill=BG_MID)
add_rect(slide, Inches(9.0), 0, Inches(0.06), H, fill=ORANGE)

# Title
add_text(slide, "TRACEON", Inches(0.6), Inches(1.2), Inches(8.0), Inches(1.8),
         size=80, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(slide, "AI-Powered Execution Flow Visualizer",
         Inches(0.6), Inches(3.0), Inches(8.0), Inches(0.7),
         size=22, color=ORANGE, align=PP_ALIGN.LEFT)
add_text(slide, "Debug smarter. Understand deeper. Build faster.",
         Inches(0.6), Inches(3.75), Inches(8.0), Inches(0.55),
         size=16, color=GRAY_LIGHT, italic=True)

# Links on right panel
add_text(slide, "Live Demo", Inches(9.3), Inches(1.5), Inches(3.7), Inches(0.45),
         size=14, bold=True, color=ORANGE)
add_text(slide, "frontend-gamma-vert-20.vercel.app",
         Inches(9.3), Inches(1.95), Inches(3.7), Inches(0.4),
         size=12, color=GRAY_LIGHT)

add_text(slide, "Stack", Inches(9.3), Inches(2.7), Inches(3.7), Inches(0.45),
         size=14, bold=True, color=ORANGE)
add_text(slide, "React · FastAPI · LLDB\nGemini AI · Vercel · MongoDB",
         Inches(9.3), Inches(3.15), Inches(3.7), Inches(0.8),
         size=12, color=GRAY_LIGHT)

add_text(slide, "Built with", Inches(9.3), Inches(4.2), Inches(3.7), Inches(0.45),
         size=14, bold=True, color=ORANGE)
add_text(slide, "React 18 + Monaco Editor\nFastAPI + python-pptx\n@react-pdf/renderer",
         Inches(9.3), Inches(4.65), Inches(3.7), Inches(1.0),
         size=12, color=GRAY_LIGHT)

add_text(slide, "nishantkumar19041@gmail.com",
         Inches(9.3), Inches(6.0), Inches(3.7), Inches(0.4),
         size=11, color=GRAY_MID)

# Bottom CTA bar
add_rect(slide, 0, H - Inches(1.0), Inches(9.0), Inches(1.0), fill=ORANGE)
add_text(slide, "Try Traceon →  frontend-gamma-vert-20.vercel.app",
         Inches(0.5), H - Inches(0.75), Inches(8.0), Inches(0.5),
         size=18, bold=True, color=WHITE)

add_rect(slide, 0, H - Inches(0.12), W, Inches(0.12), fill=ORANGE_DARK)

slide_number(slide, 13, TOTAL)


# ─── Save ────────────────────────────────────────────────────────────────────
output = "Traceon_Project.pptx"
prs.save(output)
print(f"✅  Saved → {output}  ({TOTAL} slides)")
